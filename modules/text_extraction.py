import os
import PyPDF2 
import logging
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
import pytesseract  # Added for OCR
from PIL import Image  # Added for image handling with OCR
import io  # Added for image stream handling
from typing import List, Optional, Union
from .simple_preprocess import preprocess_text

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Defined a threshold for minimum characters extracted by PyMuPDF to trigger OCR
MIN_CHARS_FOR_NON_OCR = 100


def extract_text_from_pdf(pdf_path: str) -> str:
    """Extracts text from a PDF file, handling potential OCR needs."""
    text = ""
    try:
        doc = fitz.open(pdf_path)
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            page_text = page.get_text()
            if not page_text.strip():  # If text extraction yields little, trying OCR
                logger.info(
                    f"Page {page_num+1} in {pdf_path} has little text, attempting OCR.")
                try:
                    # Use the correct method get_pixmap()
                    pix = page.get_pixmap()
                    img_bytes = pix.tobytes("png")
                    img = Image.open(io.BytesIO(img_bytes))
                    ocr_text = pytesseract.image_to_string(img)
                    text += ocr_text + "\n"
                    logger.info(f"OCR successful for page {page_num+1}.")
                except Exception as ocr_error:
                    logger.warning(
                        f"OCR failed for page {page_num+1} in {pdf_path}: {ocr_error}")
                    # Fallback: add placeholder if both failed significantly
                    if not page_text.strip():
                        text += f"[Content from page {page_num+1} could not be extracted or OCR'd]\n"

            else:
                text += page_text + "\n"
        doc.close()
    except Exception as e:
        logger.error(f"Error processing PDF {pdf_path}: {e}")
    return text


def extract_text_from_docx(file_path: str) -> str:
    """
    Extract text from DOCX files with improved structure preservation.
    Includes tables, headers, and other document elements.
    """
    try:
        doc = Document(file_path)
        content: List[str] = []

        prop = doc.core_properties
        if hasattr(prop, 'title') and prop.title:
            content.append(f"Document Title: {prop.title}")

        # Processing paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                if para.style and hasattr(para.style, 'name') and para.style.name and \
                   isinstance(para.style.name, str) and para.style.name.startswith('Heading'):
                    # Adding formatting to preserve heading structure
                    content.append(f"\n## {para.text} ##\n")
                else:
                    content.append(para.text)

        # Processing tables
        for table in doc.tables:
            table_text: List[str] = []
            for row in table.rows:
                row_text = [cell.text.strip()
                            for cell in row.cells if cell.text.strip()]
                if row_text:
                    table_text.append(" | ".join(row_text))
            if table_text:
                content.append("\nTable Content:\n" + "\n".join(table_text))

        return preprocess_text("\n".join(content))

    except Exception as e:
        logger.error(f"Failed to extract text from DOCX {file_path}: {e}")
        return f"Error extracting text from {os.path.basename(file_path)}"


def extract_text_from_ppt(ppt_path: str) -> str:
    """Extracts text from a PPT or PPTX file."""
    text = ""
    try:
        prs = Presentation(ppt_path)
        for slide in prs.slides:
            # Extract text from shapes on the slide
            for shape in slide.shapes:
                # Check if shape has text frame before accessing it
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text + " "
                        text += "\n"  # Newline after each paragraph

            # Extract text from notes slide, checking existence first
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text.strip():  # Add notes only if they contain text
                    text += "\n--- Notes ---\n"
                    text += notes_text
                    text += "\n--- End Notes ---\n"

            text += "\n--- End Slide ---\n"  # Separator between slides

    except Exception as e:
        logger.error(f"Error processing PPT/PPTX {ppt_path}: {e}")
    return text
